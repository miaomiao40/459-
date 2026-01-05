"""PPTX渲染器"""

from typing import Any, List, Dict
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from .overflow import BoundingBox, TextOverflowEngine
from .themes import ThemeColorScheme


class ShapeFactory:
    """形状工厂"""
    
    @staticmethod
    def create_text_box(slide, box: BoundingBox, text: str, font_size: float,
                        color, bold: bool = False, align=PP_ALIGN.LEFT):
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return shape
    
    @staticmethod
    def create_bullet_list(slide, box: BoundingBox, items: List[Dict],
                           font_size: float, color):
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            text = item.get('text', '') if isinstance(item, dict) else str(item)
            level = item.get('level', 0) if isinstance(item, dict) else 0
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.level = level
            p.bullet = True
        return shape
    
    @staticmethod
    def create_rectangle(slide, box: BoundingBox, color):
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape


class SlideRenderer:
    """幻灯片渲染器"""
    
    WIDTH, HEIGHT, MARGIN = 13.333, 7.5, 0.5
    
    def __init__(self, theme: ThemeColorScheme):
        self.theme = theme
        self.overflow = TextOverflowEngine()
    
    def get_layout(self, slide_type: str) -> Dict[str, BoundingBox]:
        w = self.WIDTH - 2 * self.MARGIN
        layouts = {
            'title': {'title': BoundingBox(self.MARGIN, 2.5, w, 1.5),
                     'content': BoundingBox(self.MARGIN, 4.2, w, 1.0)},
            'content': {'title': BoundingBox(self.MARGIN, 0.4, w, 0.8),
                       'content': BoundingBox(self.MARGIN, 1.4, w, 5.5)},
            'two_column': {'title': BoundingBox(self.MARGIN, 0.4, w, 0.8),
                          'left': BoundingBox(self.MARGIN, 1.4, w/2-0.2, 5.5),
                          'right': BoundingBox(self.MARGIN+w/2+0.2, 1.4, w/2-0.2, 5.5)},
            'closing': {'title': BoundingBox(self.MARGIN, 2.5, w, 1.5),
                       'content': BoundingBox(self.MARGIN, 4.2, w, 1.0)},
        }
        return layouts.get(slide_type, layouts['content'])
    
    def render(self, prs: Presentation, data: Dict) -> Any:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_type = data.get('slide_type', 'content')
        layout = self.get_layout(slide_type)
        
        if slide_type in ['title', 'closing']:
            ShapeFactory.create_rectangle(slide, BoundingBox(0, 0, self.WIDTH, 2.2),
                                         self.theme.get_rgb('primary'))
        
        title = data.get('title', '')
        if title:
            fs = 36 if slide_type in ['title', 'closing'] else 28
            result = self.overflow.fit_text(title, layout['title'], fs)
            color = self.theme.get_rgb('text_light' if slide_type in ['title', 'closing'] else 'primary')
            align = PP_ALIGN.CENTER if slide_type in ['title', 'closing'] else PP_ALIGN.LEFT
            ShapeFactory.create_text_box(slide, layout['title'], result['text'],
                                        result['font_size'], color, True, align)
        
        subtitle = data.get('subtitle')
        if subtitle and 'content' in layout:
            result = self.overflow.fit_text(subtitle, layout['content'], 24)
            ShapeFactory.create_text_box(slide, layout['content'], result['text'],
                                        result['font_size'], self.theme.get_rgb('text_dark'),
                                        False, PP_ALIGN.CENTER if slide_type == 'title' else PP_ALIGN.LEFT)
        
        points = data.get('body_points', [])
        if points and 'content' in layout:
            processed = [{'text': p.get('text', '') if isinstance(p, dict) else str(p),
                         'level': p.get('level', 0) if isinstance(p, dict) else 0}
                        for p in points]
            ShapeFactory.create_bullet_list(slide, layout['content'], processed, 18,
                                           self.theme.get_rgb('text_dark'))
        
        if slide_type == 'two_column':
            for col, key in [('left_column', 'left'), ('right_column', 'right')]:
                items = data.get(col, [])
                if items and key in layout:
                    processed = [{'text': p.get('text', ''), 'level': p.get('level', 0)}
                                if isinstance(p, dict) else {'text': str(p), 'level': 0}
                                for p in items]
                    ShapeFactory.create_bullet_list(slide, layout[key], processed, 16,
                                                   self.theme.get_rgb('text_dark'))
        
        notes = data.get('speaker_notes')
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        
        return slide


__all__ = ['ShapeFactory', 'SlideRenderer']